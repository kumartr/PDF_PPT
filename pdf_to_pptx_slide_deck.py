
import fitz  # PyMuPDF
import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def pdf_to_images(pdf_path, image_folder):
    os.makedirs(image_folder, exist_ok=True)
    doc = fitz.open(pdf_path)
    image_paths = []

    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # High resolution
        image_path = os.path.join(image_folder, f"page_{i+1}.png")
        pix.save(image_path)
        image_paths.append(image_path)

    print(f"✅ Converted {len(image_paths)} pages to images.")
    return image_paths

def images_to_pptx(image_paths, pptx_path):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout

    for image_path in image_paths:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width)

    prs.save(pptx_path)
    print(f"✅ PowerPoint saved to '{pptx_path}'")

# Example usage
pdf_file = "sample_research_paper.pdf"  # Replace with your PDF
output_folder = "pdf_page_images"
pptx_output = "pdf_to_slides.pptx"

image_files = pdf_to_images(pdf_file, output_folder)
images_to_pptx(image_files, pptx_output)
